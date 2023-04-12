#!/usr/bin/env python

import sys
import collections
import collections.abc
from pptx import Presentation

def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    all_text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text)

    return all_text

def main():
    if len(sys.argv) < 2:
        print("Please provide the PowerPoint file path as an argument.")
        sys.exit(1)

    file_path = sys.argv[1]
    all_text = extract_text_from_ppt(file_path)

    with open("ppt_text.txt", "w") as file:
        for index, slide_text in enumerate(all_text):
            file.write(f"Slide {index + 1}:\n{slide_text}\n")

print("The extracted text has been saved to ppt_text.txt")

if __name__ == "__main__":
    main()